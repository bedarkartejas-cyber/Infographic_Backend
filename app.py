import os
import re
import json
import time
import uuid
from io import BytesIO
from typing import Optional
import asyncio
from datetime import datetime

from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Header
from fastapi.responses import StreamingResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from pptx import Presentation
import requests

from supabase_db import db
from main import generate_marketing_assets_stream_supabase, generate_marketing_assets_supabase

app = FastAPI(
    title="Marketing Generator API",
    description="Generate marketing assets from PPTX or website content",
    version="2.0.0"
)

# Serve local images if needed
os.makedirs("local_images", exist_ok=True)
app.mount("/local_images", StaticFiles(directory="local_images"), name="local_images")

# --- CORS Configuration ---
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- CONFIGURATION ---
TIMEOUT = (10, 60)
HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
        "AppleWebKit/537.36 (KHTML, like Gecko) "
        "Chrome/120.0 Safari/537.36"
    )
}

def extract_text_from_pptx(file_content):
    """Extract text from PPTX file"""
    try:
        prs = Presentation(BytesIO(file_content))
        full_text_output = []

        for i, slide in enumerate(prs.slides, start=1):
            slide_content = []
            slide_content.append(f"--- SLIDE {i} ---")

            try:
                if slide.shapes.title and slide.shapes.title.text.strip():
                    title_text = slide.shapes.title.text.strip()
                    slide_content.append(f"[Title]: {title_text}")
            except:
                pass

            text_shapes = []
            for shape in slide.shapes:
                if not shape.has_text_frame: 
                    continue
                if shape == slide.shapes.title: 
                    continue
                text_shapes.append(shape)
            
            text_shapes.sort(key=lambda s: (s.top, s.left))

            for shape in text_shapes:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_content.append(text)

            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame:
                    notes_text = notes_frame.text.strip()
                    if notes_text:
                        slide_content.append(f"\n[Speaker Notes]:\n{notes_text}")

            full_text_output.append("\n".join(slide_content))

        return "\n\n".join(full_text_output)

    except Exception as e:
        print(f"Error reading PPTX: {e}")
        return ""

def extract_text_from_url_sync(url: str):
    """Sync web scraping"""
    try:
        response = requests.get(url, headers=HEADERS, timeout=TIMEOUT)
        response.raise_for_status()
        html_content = response.text
        
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            return ""
        
        soup = BeautifulSoup(html_content, "html.parser")
        for tag in soup(["script", "style", "noscript", "header", "footer", "nav", "aside"]):
            tag.decompose()

        text = soup.get_text(separator="\n")
        text = re.sub(r"\n\s*\n+", "\n\n", text)
        text = re.sub(r"[ \t]+", " ", text)
        return text.strip()
        
    except Exception as e:
        print(f"Error scraping URL: {e}")
        return ""

@app.get("/")
async def index():
    return JSONResponse(
        content={
            "message": "🚀 Marketing Generator API",
            "version": "2.0.0",
            "status": "running",
            "features": [
                "PPTX content extraction",
                "Website content scraping", 
                "Marketing brief generation",
                "Creative angles generation",
                "Email copy generation",
                "Image prompt generation",
                "Parallel image generation (A2E)",
                "Supabase storage integration"
            ],
            "endpoints": {
                "POST /api/generate": "Generate all assets at once",
                "POST /api/generate-stream": "Stream generation progress",
                "GET /api/generations": "Get user's generations",
                "GET /api/generations/{id}": "Get specific generation",
                "GET /api/status/{id}": "Check generation status"
            },
            "storage": "Supabase + Local fallback"
        }
    )

@app.get("/health")
async def health_check():
    """Health check endpoint"""
    return {
        "status": "healthy",
        "timestamp": time.time(),
        "datetime": datetime.now().isoformat(),
        "service": "marketing-generator"
    }

@app.get("/api/test")
async def test_endpoint():
    """Test endpoint for debugging"""
    return {
        "status": "ok",
        "supabase_connected": db._get_client() is not None,
        "memory_generations": len(db._memory_generations),
        "memory_images": len(db._memory_images)
    }

async def get_user_id(x_user_id: Optional[str] = Header(None)):
    """Get user ID from header or use default"""
    if x_user_id:
        return x_user_id
    return "demo-user"

@app.post("/api/generate")
async def generate_api(
    website_url: Optional[str] = Form(None),
    ppt_file: Optional[UploadFile] = File(None),
    image_count: int = Form(3),
    x_user_id: Optional[str] = Header(None)
):
    """
    Generate all marketing assets at once
    
    - website_url: Optional website URL to scrape content from
    - ppt_file: Optional PPTX file upload
    - image_count: Number of images to generate (1-5)
    - x_user_id: User identifier (optional)
    """
    start_time = time.time()
    user_id = x_user_id or "demo-user"
    
    print(f"\n🎯 GENERATION REQUEST")
    print(f"   User: {user_id}")
    print(f"   Website: {website_url}")
    print(f"   PPT File: {ppt_file.filename if ppt_file else 'None'}")
    print(f"   Images: {image_count}")
    
    website_text = ""
    ppt_text = ""
    
    try:
        # Extract website content
        if website_url:
            print(f"   Scraping website: {website_url}")
            website_text = await asyncio.to_thread(extract_text_from_url_sync, website_url)
            print(f"   Website text length: {len(website_text)} chars")
        
        # Extract PPT content
        if ppt_file and ppt_file.filename and ppt_file.filename.endswith(('.pptx', '.ppt')):
            print(f"   Processing PPT file: {ppt_file.filename}")
            file_content = await ppt_file.read()
            ppt_text = extract_text_from_pptx(file_content)
            print(f"   PPT text length: {len(ppt_text)} chars")
        
        if not ppt_text and not website_text:
            raise HTTPException(
                status_code=400, 
                detail="No content found. Please provide a valid PPTX file or Website URL."
            )
        
        # Create generation session
        generation_id = db.create_generation_session(
            user_id=user_id,
            website_url=website_url,
            ppt_text=ppt_text,
            website_text=website_text
        )
        
        print(f"   Generation ID: {generation_id}")
        
        try:
            # Limit image count for performance
            actual_image_count = min(max(1, image_count), 5)
            if actual_image_count != image_count:
                print(f"   ⚠️  Adjusted image count to {actual_image_count} (max 5)")
            
            # Run generation
            results = await generate_marketing_assets_supabase(
                ppt_text=ppt_text,
                website_text=website_text,
                user_id=user_id,
                generation_id=generation_id,
                image_count=actual_image_count
            )
            
            total_time = time.time() - start_time
            results["generation_time"] = round(total_time, 2)
            
            print(f"✅ Generation completed in {total_time:.1f}s")
            print(f"   Images generated: {len(results.get('generated_images', []))}")
            
            return {
                "success": True,
                "generation_id": generation_id,
                "user_id": user_id,
                "data": results,
                "performance": {
                    "total_time": round(total_time, 2),
                    "images_generated": len(results.get('generated_images', [])),
                    "storage": results.get('performance', {}).get('storage', 'unknown')
                }
            }
            
        except Exception as gen_error:
            # Mark generation as failed
            print(f"❌ Generation failed: {gen_error}")
            db.fail_generation(generation_id, str(gen_error))
            raise HTTPException(
                status_code=500, 
                detail=f"Generation failed: {str(gen_error)}"
            )
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ Unexpected error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/generate-stream")
async def generate_stream_api(
    website_url: Optional[str] = Form(None),
    ppt_file: Optional[UploadFile] = File(None),
    image_count: int = Form(3),
    x_user_id: Optional[str] = Header(None)
):
    """
    Stream generation progress in real-time
    
    Returns NDJSON stream with events:
    - start: Generation started
    - brief: Marketing brief ready
    - email: Email copy ready  
    - image_start: Image generation starting
    - image: Individual image ready
    - complete: All done
    - error: Something went wrong
    """
    user_id = x_user_id or "demo-user"
    
    print(f"\n🌀 STREAMING GENERATION REQUEST")
    print(f"   User: {user_id}")
    
    try:
        website_text = ""
        ppt_text = ""
        
        # Extract content
        if website_url:
            website_text = await asyncio.to_thread(extract_text_from_url_sync, website_url)
        
        if ppt_file and ppt_file.filename and ppt_file.filename.endswith(('.pptx', '.ppt')):
            file_content = await ppt_file.read()
            ppt_text = extract_text_from_pptx(file_content)
        
        if not ppt_text and not website_text:
            raise HTTPException(status_code=400, detail="No content found.")
        
        # Create generation session
        generation_id = db.create_generation_session(
            user_id=user_id,
            website_url=website_url,
            ppt_text=ppt_text,
            website_text=website_text
        )
        
        # Limit image count
        actual_image_count = min(max(1, image_count), 5)
        
        async def generate():
            try:
                # Send start event
                yield json.dumps({
                    "type": "start",
                    "timestamp": time.time(),
                    "generation_id": generation_id,
                    "user_id": user_id,
                    "message": "Starting marketing generation...",
                    "image_count": actual_image_count
                }) + "\n"
                
                # Stream generation
                async for chunk in generate_marketing_assets_stream_supabase(
                    ppt_text=ppt_text,
                    website_text=website_text,
                    user_id=user_id,
                    generation_id=generation_id,
                    image_count=actual_image_count
                ):
                    yield json.dumps(chunk) + "\n"
                
                # Send completion
                yield json.dumps({
                    "type": "complete",
                    "timestamp": time.time(),
                    "generation_id": generation_id,
                    "message": "Generation completed successfully!"
                }) + "\n"
                
            except Exception as e:
                # Mark as failed
                db.fail_generation(generation_id, str(e))
                
                # Send error
                yield json.dumps({
                    "type": "error",
                    "timestamp": time.time(),
                    "message": str(e),
                    "generation_id": generation_id
                }) + "\n"
        
        return StreamingResponse(
            generate(),
            media_type='application/x-ndjson',
            headers={
                "X-Accel-Buffering": "no",
                "Cache-Control": "no-cache",
                "X-Generation-ID": generation_id,
                "X-User-ID": user_id
            }
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/generations")
async def get_user_generations(
    x_user_id: Optional[str] = Header(None),
    limit: int = 20
):
    """Get all generations for current user"""
    try:
        user_id = x_user_id or "demo-user"
        
        # For memory storage
        user_generations = []
        for gen_id, gen_data in db._memory_generations.items():
            if gen_data.get("user_id") == user_id:
                # Get images for this generation
                images = db._memory_images.get(gen_id, [])
                gen_data["images"] = images
                user_generations.append(gen_data)
        
        # Sort by created_at (newest first)
        user_generations.sort(
            key=lambda x: x.get("created_at", ""), 
            reverse=True
        )
        
        # Apply limit
        user_generations = user_generations[:limit]
        
        return {
            "success": True,
            "user_id": user_id,
            "count": len(user_generations),
            "generations": user_generations,
            "storage": "memory"
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/generations/{generation_id}")
async def get_generation(
    generation_id: str,
    x_user_id: Optional[str] = Header(None)
):
    """Get specific generation by ID"""
    try:
        user_id = x_user_id or "demo-user"
        
        generation = db.get_generation(generation_id, user_id)
        
        if not generation:
            raise HTTPException(
                status_code=404, 
                detail=f"Generation {generation_id} not found or access denied"
            )
        
        return {
            "success": True,
            "generation": generation,
            "storage": generation.get("storage", "unknown")
        }
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/status/{generation_id}")
async def get_generation_status(
    generation_id: str,
    x_user_id: Optional[str] = Header(None)
):
    """Check status of a generation"""
    try:
        user_id = x_user_id or "demo-user"
        
        generation = db.get_generation(generation_id, user_id)
        
        if not generation:
            raise HTTPException(status_code=404, detail="Generation not found")
        
        status_info = {
            "generation_id": generation_id,
            "status": generation.get("status", "unknown"),
            "user_id": generation.get("user_id"),
            "created_at": generation.get("created_at"),
            "updated_at": generation.get("updated_at"),
            "total_images": generation.get("total_images", 0),
            "completed_images": generation.get("completed_images", 0),
            "images_count": len(generation.get("images", [])),
            "storage": generation.get("storage", "unknown")
        }
        
        return {
            "success": True,
            "status": status_info
        }
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/generations/{generation_id}")
async def delete_generation(
    generation_id: str,
    x_user_id: Optional[str] = Header(None)
):
    """Delete a generation and all its images"""
    try:
        user_id = x_user_id or "demo-user"
        
        # Check if generation exists and belongs to user
        generation = db.get_generation(generation_id, user_id)
        if not generation:
            raise HTTPException(status_code=404, detail="Generation not found")
        
        # Delete from memory storage
        if generation_id in db._memory_generations:
            del db._memory_generations[generation_id]
        
        if generation_id in db._memory_images:
            del db._memory_images[generation_id]
        
        # TODO: Implement Supabase deletion when using database
        
        return {
            "success": True,
            "message": f"Generation {generation_id} deleted",
            "generation_id": generation_id
        }
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/system-info")
async def system_info():
    """Get system information and configuration"""
    import psutil
    
    cpu_percent = psutil.cpu_percent(interval=1)
    memory = psutil.virtual_memory()
    
    return {
        "system": {
            "cpu_usage_percent": cpu_percent,
            "memory_usage_percent": memory.percent,
            "memory_available_gb": round(memory.available / (1024**3), 2),
            "python_version": os.sys.version
        },
        "api": {
            "max_images": 5,
            "parallel_generation": True,
            "storage_backend": "Supabase + Local fallback"
        },
        "tables": {
            "marketing_generations": "Stores generation sessions",
            "marketing_images": "Stores generated images"
        }
    }

# Error handlers
@app.exception_handler(HTTPException)
async def http_exception_handler(request, exc):
    return JSONResponse(
        status_code=exc.status_code,
        content={
            "success": False,
            "error": exc.detail,
            "status_code": exc.status_code
        }
    )

@app.exception_handler(Exception)
async def general_exception_handler(request, exc):
    return JSONResponse(
        status_code=500,
        content={
            "success": False,
            "error": str(exc),
            "type": type(exc).__name__
        }
    )

if __name__ == "__main__":
    import uvicorn
    
    print("\n" + "="*60)
    print("🚀 MARKETING GENERATOR API")
    print("="*60)
    print("Starting server...")
    print(f"Time: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    print("="*60)
    
    # Create required directories
    os.makedirs("local_images", exist_ok=True)
    
    uvicorn.run(
        app, 
        host="0.0.0.0", 
        port=5000, 
        reload=True,
        log_level="info",
        timeout_keep_alive=300
    )