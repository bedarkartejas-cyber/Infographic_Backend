import re
import requests
from io import BytesIO
from pptx import Presentation
import logging

logger = logging.getLogger(__name__)

# Configuration
TIMEOUT = (10, 60)
HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
        "AppleWebKit/537.36 (KHTML, like Gecko) "
        "Chrome/120.0 Safari/537.36"
    )
}

def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from PPTX file"""
    try:
        prs = Presentation(BytesIO(file_content))
        full_text_output = []

        for i, slide in enumerate(prs.slides, start=1):
            slide_content = []
            slide_content.append(f"--- SLIDE {i} ---")

            try:
                if slide.shapes.title and slide.shapes.title.text.strip():
                    title_text = slide.shapes.title.text.strip()
                    slide_content.append(f"[Title]: {title_text}")
            except:
                pass

            text_shapes = []
            for shape in slide.shapes:
                if not shape.has_text_frame: 
                    continue
                if shape == slide.shapes.title: 
                    continue
                text_shapes.append(shape)
            
            text_shapes.sort(key=lambda s: (s.top, s.left))

            for shape in text_shapes:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_content.append(text)

            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame:
                    notes_text = notes_frame.text.strip()
                    if notes_text:
                        slide_content.append(f"\n[Speaker Notes]:\n{notes_text}")

            full_text_output.append("\n".join(slide_content))

        result = "\n\n".join(full_text_output)
        logger.info(f"Extracted {len(result)} chars from PPTX")
        return result

    except Exception as e:
        logger.error(f"Error reading PPTX: {e}")
        return ""

def extract_text_from_url_sync(url: str, timeout: int = 30) -> str:
    """Sync web scraping with configurable timeout"""
    try:
        response = requests.get(url, headers=HEADERS, timeout=timeout)
        response.raise_for_status()
        html_content = response.text
        
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            logger.error("BeautifulSoup not installed")
            return ""
        
        soup = BeautifulSoup(html_content, "html.parser")
        for tag in soup(["script", "style", "noscript", "header", "footer", "nav", "aside"]):
            tag.decompose()

        text = soup.get_text(separator="\n")
        text = re.sub(r"\n\s*\n+", "\n\n", text)
        text = re.sub(r"[ \t]+", " ", text)
        
        result = text.strip()
        logger.info(f"Extracted {len(result)} chars from URL: {url}")
        return result
        
    except requests.exceptions.Timeout:
        logger.error(f"Timeout scraping URL: {url}")
        return ""
    except requests.exceptions.RequestException as e:
        logger.error(f"Request error scraping URL {url}: {e}")
        return ""
    except Exception as e:
        logger.error(f"Unexpected error scraping URL {url}: {e}")
        return ""